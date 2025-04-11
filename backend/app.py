from flask import Flask, request, jsonify
from flask_cors import CORS
import os
import fitz
import re
import time
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
import json
import requests
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np
import pickle
import docx
from pptx import Presentation
import hashlib
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import threading

# Initialize Flask app
app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

# Constants
UPLOAD_FOLDER = 'uploads'
MODEL_DIR = 'models'
INDEX_PATH = os.path.join(MODEL_DIR, 'faiss_index.pkl')
METADATA_PATH = os.path.join(MODEL_DIR, 'metadata.pkl')

# TODO: hetha document eli ysiir alih imbedding automatik : Badlou bel forlder eli al serveur ydir alih imbedding
DOCUMENTS_FOLDER = 'uploads'  # Folder to monitor for automatic indexing 
PROCESSED_FILES_REGISTRY = os.path.join(MODEL_DIR, 'processed_files.pkl')

# Create necessary directories
for directory in [UPLOAD_FOLDER, MODEL_DIR]:
    if not os.path.exists(directory):
        os.makedirs(directory)

# Create the documents folder if it doesn't exist
if not os.path.exists(DOCUMENTS_FOLDER):
    os.makedirs(DOCUMENTS_FOLDER)

# File tracking system
processed_files = {}    

# Initialize the embedding model
embedding_model = SentenceTransformer('sentence-transformers/all-MiniLM-L6-v2')
embedding_dimension = 384  # Dimension for the all-MiniLM-L6-v2 model

# FAISS index and metadata storage
index = None
metadata = []

# Initialize or load FAISS index
def initialize_index():
    global index, metadata
    if os.path.exists(INDEX_PATH) and os.path.exists(METADATA_PATH):
        # Load existing index and metadata
        with open(INDEX_PATH, 'rb') as f:
            index = pickle.load(f)
        with open(METADATA_PATH, 'rb') as f:
            metadata = pickle.load(f)
        print(f"Loaded existing index with {index.ntotal} vectors")
    else:
        # Create a new index
        index = faiss.IndexFlatL2(embedding_dimension)
        metadata = []
        print("Created new FAISS index")

# Call initialization
initialize_index()

# Routes
@app.route('/', methods=['GET'])
def home():
    """Root endpoint for testing."""
    return jsonify({'message': 'Welcome to the Local RAG Chatbot API!'})
@app.route('/api/upload', methods=['POST'])
def upload_file():
    """Upload file endpoint for multiple document formats (PDF, DOCX, PPTX)."""
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    # Check file extension
    filename = secure_filename(file.filename)
    file_extension = os.path.splitext(filename)[1].lower()
    
    # Validate file extension
    allowed_extensions = ['.pdf', '.docx', '.pptx']
    if file_extension not in allowed_extensions:
        return jsonify({'error': f'Unsupported file format. Allowed formats: {", ".join(allowed_extensions)}'}), 400

    try:
        # Save file in both directories
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        print(f"File saved successfully at {filepath}")

        # # Save in static directory for permanent storage
        # static_filepath = os.path.join(STATIC_DIR, filename)
        # with open(static_filepath, 'wb') as f_static:
        #     file.seek(0)  # Reset file pointer to beginning
        #     f_static.write(file.read())
        # print(f"File saved permanently at {static_filepath}")

        # Extract text based on file type
        if file_extension == '.pdf':
            text = extract_text_from_pdf(filepath)
        elif file_extension == '.docx':
            text = extract_text_from_docx(filepath)
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(filepath)
        if text:
            print(f"Text extracted successfully from {file_extension} file. Text length: {len(text)}")
        
        # Chunk the text
        chunks = chunk_text_simple(text, chunk_size=1000, overlap=100)
        chunk_count = len(chunks)
        print(f"Text split into {chunk_count} chunks.")
        
        # Store chunks in FAISS index
        success_count = store_in_faiss(chunks, filename)
        
        print(f"Successfully stored {success_count} out of {chunk_count} chunks in FAISS index.")
        return jsonify({
            'message': f'File processed with {success_count} chunks stored',
            'filename': filename,
            'format': file_extension[1:],  # Remove the leading dot
            'chunks': chunk_count
        })
        
    except Exception as e:
        print(f"Error in upload_file: {str(e)}")
        return jsonify({'error': f'Error processing file: {str(e)}'}), 500

@app.route('/api/chat', methods=['POST'])
def chat_endpoint():
    """Chat endpoint for handling user queries."""
    data = request.json
    query = data.get('query')
    history = data.get('history', None)  # This should be an array of context tokens
    
    if not query:
        return jsonify({'error': 'No query provided'}), 400
    
    try:
        # Retrieve context using FAISS
        augmented_query, contexts = augment_prompt(query)
        print(f"Augmented query with context: {augmented_query[:100]}...")
        
        # Generate response using Ollama
        response = generate_response_ollama(augmented_query, history)
        
        if not response:
            return jsonify({'error': 'Failed to generate response'}), 500
        
        return jsonify({
            'response': response,
            'augmentedQuery': augmented_query,
            'contexts': contexts
        })
    except Exception as e:
        print(f"Error in chat endpoint: {e}")
        return jsonify({'error': f'Error processing query: {str(e)}'}), 500
@app.route('/api/test', methods=['GET'])
def test_endpoint():
    """Test endpoint to verify API is working."""
    return jsonify({'message': 'Hello from the Local RAG Chatbot API!'})


# Add function to manually reindex all documents in the folder
@app.route('/api/reindex-all', methods=['POST'])
def reindex_all_documents():
    """Reindex all documents in the documents folder."""
    try:
        # Reset the processed files registry to force reprocessing
        global processed_files
        processed_files = {}
        
        # Get all document files
        files_processed = 0
        for filename in os.listdir(DOCUMENTS_FOLDER):
            filepath = os.path.join(DOCUMENTS_FOLDER, filename)
            if os.path.isfile(filepath):
                file_extension = os.path.splitext(filepath)[1].lower()
                if file_extension in ['.pdf', '.docx', '.pptx']:
                    try:
                        process_document_from_path(filepath)
                        files_processed += 1
                    except Exception as e:
                        print(f"Error processing {filepath}: {e}")
        
        # Save the updated registry
        save_processed_files()
        
        return jsonify({
            'message': f'Reindexed {files_processed} documents from {DOCUMENTS_FOLDER}',
            'files_processed': files_processed
        })
    except Exception as e:
        print(f"Error reindexing documents: {e}")
        return jsonify({'error': f'Error reindexing documents: {str(e)}'}), 500

# Add endpoint to list all indexed documents
@app.route('/api/list-documents', methods=['GET'])
def list_indexed_documents():
    """List all indexed documents."""
    try:
        # Extract unique document sources from metadata
        document_sources = set()
        for item in metadata:
            if 'source' in item:
                document_sources.add(item['source'])
        
        # Also list documents in the documents folder
        folder_documents = []
        for filename in os.listdir(DOCUMENTS_FOLDER):
            filepath = os.path.join(DOCUMENTS_FOLDER, filename)
            if os.path.isfile(filepath):
                file_extension = os.path.splitext(filepath)[1].lower()
                if file_extension in ['.pdf', '.docx', '.pptx']:
                    folder_documents.append({
                        'filename': filename,
                        'path': filepath,
                        'indexed': filename in document_sources,
                        'size': os.path.getsize(filepath)
                    })
        
        return jsonify({
            'indexed_documents': list(document_sources),
            'folder_documents': folder_documents
        })
    except Exception as e:
        print(f"Error listing documents: {e}")
        return jsonify({'error': f'Error listing documents: {str(e)}'}), 500

# Helper Functions

# Load processed files registry if exists
def load_processed_files():
    global processed_files
    if os.path.exists(PROCESSED_FILES_REGISTRY):
        with open(PROCESSED_FILES_REGISTRY, 'rb') as f:
            processed_files = pickle.load(f)
            print(f"Loaded registry with {len(processed_files)} processed files")
    else:
        processed_files = {}
        print("Created new processed files registry")

# Save processed files registry
def save_processed_files():
    with open(PROCESSED_FILES_REGISTRY, 'wb') as f:
        pickle.dump(processed_files, f)
    print(f"Saved registry with {len(processed_files)} processed files")

# Check if a file has been processed before
def is_file_processed(filepath):
    # Calculate file hash for reliable tracking
    try:
        with open(filepath, 'rb') as f:
            file_hash = hashlib.md5(f.read()).hexdigest()
        
        # Check if file hash exists in registry
        if filepath in processed_files and processed_files[filepath] == file_hash:
            return True
        
        # Update registry with new hash
        processed_files[filepath] = file_hash
        save_processed_files()
        return False
    except Exception as e:
        print(f"Error checking if file is processed: {e}")
        return False

# File event handler for watchdog
class DocumentHandler(FileSystemEventHandler):
    def on_created(self, event):
        self.process_file_event(event)
    
    def on_modified(self, event):
        self.process_file_event(event)
        
    def process_file_event(self, event):
        # Skip directory events and hidden files
        if event.is_directory or os.path.basename(event.src_path).startswith('.'):
            return
            
        # Check file extension
        file_extension = os.path.splitext(event.src_path)[1].lower()
        allowed_extensions = ['.pdf', '.docx', '.pptx']
        
        if file_extension in allowed_extensions:
            # Add a small delay to ensure the file is fully written
            time.sleep(1)
            
            # Check if we've already processed this file version
            if not is_file_processed(event.src_path):
                print(f"New or modified document detected: {event.src_path}")
                try:
                    # Process the document
                    process_document_from_path(event.src_path)
                except Exception as e:
                    print(f"Error processing document {event.src_path}: {e}")

# Process a document from the filesystem
def process_document_from_path(filepath):
    """Process a document from the filesystem and add to index."""
    try:
        filename = os.path.basename(filepath)
        print(f"Processing document: {filename}")
        
        # Extract text based on file type
        file_extension = os.path.splitext(filepath)[1].lower()
        
        if file_extension == '.pdf':
            text = extract_text_from_pdf(filepath)
        elif file_extension == '.docx':
            text = extract_text_from_docx(filepath)
        elif file_extension == '.pptx':
            text = extract_text_from_pptx(filepath)
        else:
            print(f"Unsupported file format: {file_extension}")
            return
        
        print(f"Text extracted successfully from {filepath}. Text length: {len(text)}")
        
        # Chunk the text
        chunks = chunk_text_simple(text, chunk_size=1000, overlap=100)
        chunk_count = len(chunks)
        print(f"Text split into {chunk_count} chunks.")
        
        # Store chunks in FAISS index
        success_count = store_in_faiss(chunks, filename)
        
        print(f"Successfully stored {success_count} out of {chunk_count} chunks in FAISS index.")
        return success_count
        
    except Exception as e:
        print(f"Error processing document {filepath}: {e}")
        raise e

# Start the file system observer
def start_document_observer():
    event_handler = DocumentHandler()
    observer = Observer()
    observer.schedule(event_handler, DOCUMENTS_FOLDER, recursive=False)
    observer.start()
    print(f"Started monitoring {DOCUMENTS_FOLDER} for new documents")
    return observer

def extract_text_from_pdf(filepath):
    """Extract text from a PDF file."""
    text = ""
    try:
        doc = fitz.open(filepath)
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        raise Exception(f"Failed to extract text from PDF: {e}")

def extract_text_from_docx(filepath):
    """Extract text from a Word document."""
    text = ""
    try:
        doc = docx.Document(filepath)
        for para in doc.paragraphs:
            text += para.text + "\n"
        
        # Extract text from tables as well
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return text
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        raise Exception(f"Failed to extract text from DOCX: {e}")

def extract_text_from_pptx(filepath):
    """Extract text from a PowerPoint presentation."""
    text = ""
    try:
        pres = Presentation(filepath)
        for slide in pres.slides:
            text += f"Slide {pres.slides.index(slide) + 1}:\n"
            
            # Extract text from shapes (includes text boxes)
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            
            # Add separator between slides
            text += "\n" + "-"*40 + "\n"
        
        return text
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        raise Exception(f"Failed to extract text from PPTX: {e}")

def process_document(filepath):
    """Process document based on file extension."""
    file_extension = os.path.splitext(filepath)[1].lower()
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(filepath)
    elif file_extension == '.docx':
        return extract_text_from_docx(filepath)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(filepath)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")


def chunk_text_simple(text, chunk_size=1000, overlap=100):
    """Split text into overlapping chunks with a simpler method."""
    chunks = []
    start = 0
    text_length = len(text)
    
    # If text is very short, return it as is
    if text_length <= chunk_size:
        return [text]
    
    while start < text_length:
        # Determine the end of current chunk
        end = min(start + chunk_size, text_length)
        
        # Add chunk to results
        current_chunk = text[start:end]
        chunks.append(current_chunk)
        
        # Calculate start of next chunk with overlap
        start = start + chunk_size - overlap
    
    return chunks

def store_in_faiss(chunks, filename):
    """Store processed chunks in FAISS index."""
    global index, metadata
    
    success_count = 0
    for i, chunk in enumerate(chunks):
        try:
            # Generate embedding
            embedding = embedding_model.encode([chunk])[0]
            
            # Convert to the format FAISS needs
            embedding_np = np.array([embedding]).astype('float32')
            
            # Add to FAISS index
            index.add(embedding_np)
            
            # Store metadata
            chunk_id = f"{filename}-chunk-{i}"
            # Truncate text for metadata
            truncated_text = chunk[:500] if len(chunk) > 500 else chunk
            # Clean text of non-ASCII characters
            clean_text = ''.join(char for char in truncated_text if ord(char) < 128)
            
            metadata.append({
                'text': clean_text, 
                'full_text': chunk,
                'source': filename, 
                'chunk_id': chunk_id
            })
            
            success_count += 1
            
        except Exception as chunk_error:
            print(f"Error storing chunk {i}: {chunk_error}")
    
    # Save updated index and metadata
    with open(INDEX_PATH, 'wb') as f:
        pickle.dump(index, f)
    with open(METADATA_PATH, 'wb') as f:
        pickle.dump(metadata, f)
        
    return success_count

def augment_prompt(query, top_k=3):
    """Augment user query with context from FAISS."""
    try:
        # Generate embedding for query
        query_embedding = embedding_model.encode([query])[0]
        query_embedding_np = np.array([query_embedding]).astype('float32')
        
        # Search FAISS index
        D, I = index.search(query_embedding_np, top_k)
        
        # Process results
        contexts = []
        for i, idx in enumerate(I[0]):
            if idx < len(metadata) and idx >= 0:  # Check if index is valid
                doc = metadata[idx]
                contexts.append({
                    'content': doc['full_text'],
                    'source': doc['source'],
                    'chunk_id': doc['chunk_id'],
                    'similarity': float(D[0][i])
                })
        
        # Check if contexts were found
        if not contexts:
            print("No relevant contexts found in FAISS.")
            return f"Réponds à cette requête: {query}", []
        
        # Build context from results
        context_texts = []
        for i, doc in enumerate(contexts):
            context_texts.append(f"Document {i+1} (Source: {doc['source']}):\n{doc['content']}")
        
        source_knowledge = "\n\n".join(context_texts)
        print(f"Found {len(contexts)} relevant contexts for query.")
        
        augmented_query = f"""En utilisant UNIQUEMENT les contextes ci-dessous, réponds à la requête. 
Si les contextes ne contiennent pas d'information pertinente pour répondre à la requête, indique que tu n'as pas suffisamment d'informations.

Contextes:
{source_knowledge}

Requête: {query}"""
        
        return augmented_query, contexts
    except Exception as e:
        print(f"Error in augment_prompt: {e}")
        # Fallback to regular query without context
        return f"Réponds à cette requête: {query}", []
import requests

def generate_response_ollama(augmented_query, history=None):
    """Generate response using Ollama API."""
    try:
        # Prepare the prompt with system message
        system_message = "Vous êtes un assistant utile qui répond aux questions basées sur les contextes fournis. Si les contextes ne contiennent pas l'information demandée, indiquez clairement que vous ne pouvez pas répondre à la question avec les données dont vous disposez."
        
        # Construct the payload for Ollama
        payload = {
            "model": "deepseek-r1:7b",
            "prompt": augmented_query,
            "system": system_message,
            "stream": False
        }
        
        # Add chat history if provided
        if history and isinstance(history, list):
            payload["context"] = history
            print(payload)

        # Make the API call to Ollama
        response = requests.post(
            "http://localhost:11434/api/generate",
            json=payload
        )
        
        # Check if the request was successful
        if response.status_code == 200:
            result = response.json()
            response_content = result.get("response", "")

            # Supprimer les balises <think> et </think> avec split() puis enlever les espaces inutiles
            if "<think>" in response_content and "</think>" in response_content:
                response_content = response_content.split("<think>")[0] + response_content.split("</think>")[-1]

            # Supprimer les espaces avant et après la réponse
            response_content = response_content.strip()

            print(f"Response generated successfully.")
            return response_content
        else:
            print(f"Error from Ollama API: {response.status_code}, {response.text}")
            return "Désolé, je n'ai pas pu générer une réponse. Veuillez réessayer."
            
    except Exception as e:
        print(f"Error generating response: {e}")
        return "Désolé, je n'ai pas pu générer une réponse. Veuillez réessayer."


# Add this to your main function
if __name__ == '__main__':
    # Load existing processed files registry
    load_processed_files()
    
    # Start the document observer in a background thread
    observer = start_document_observer()
    
    try:
        # Initial indexing of documents in the folder at startup
        print("Performing initial indexing of documents folder...")
        for filename in os.listdir(DOCUMENTS_FOLDER):
            filepath = os.path.join(DOCUMENTS_FOLDER, filename)
            if os.path.isfile(filepath) and not filename.startswith('.'):
                file_extension = os.path.splitext(filepath)[1].lower()
                if file_extension in ['.pdf', '.docx', '.pptx']:
                    if not is_file_processed(filepath):
                        try:
                            process_document_from_path(filepath)
                        except Exception as e:
                            print(f"Error during initial indexing of {filepath}: {e}")
        
        # Start Flask app
        print("Starting Flask application...")
        app.run(debug=True, host='0.0.0.0', port=int(os.environ.get("PORT", 5001)))
    
    finally:
        # Stop the observer when the app is shutting down
        observer.stop()
        observer.join()
        print("Document monitoring stopped")